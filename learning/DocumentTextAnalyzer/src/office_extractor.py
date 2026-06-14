"""
office_extractor.py
Extrahiert Text aus Office-Dateien: DOCX, XLSX, PPTX.
"""

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation


def extract_docx(file_path: str) -> str:
    """
    Liest Text aus einer Word-Datei (.docx).
    Gibt alle Absätze als Text zurück.
    """
    try:
        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

        if not paragraphs:
            return "Kein Text in der Word-Datei gefunden."

        return "\n".join(paragraphs)

    except Exception as e:
        return f"Fehler beim Lesen der Word-Datei: {e}"


def extract_xlsx(file_path: str) -> str:
    """
    Liest Zellinhalte aus einer Excel-Datei (.xlsx).
    Gibt alle Blätter und deren Zellinhalte als Text zurück.
    """
    try:
        wb = load_workbook(file_path, read_only=True, data_only=True)
        text_parts = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            text_parts.append(f"=== Tabellenblatt: {sheet_name} ===")

            for row in ws.iter_rows(values_only=True):
                row_values = [str(cell) for cell in row if cell is not None]
                if row_values:
                    text_parts.append("\t".join(row_values))

        num_sheets = len(wb.sheetnames)
        wb.close()

        if len(text_parts) <= num_sheets:
            return "Keine Inhalte in der Excel-Datei gefunden."

        return "\n".join(text_parts)

    except Exception as e:
        return f"Fehler beim Lesen der Excel-Datei: {e}"


def extract_pptx(file_path: str) -> str:
    """
    Liest Text aus einer PowerPoint-Datei (.pptx).
    Gibt alle Folien und deren Textbausteine zurück.
    """
    try:
        prs = Presentation(file_path)
        text_parts = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            slide_texts.append(line)

            if slide_texts:
                text_parts.append(f"--- Folie {slide_num} ---")
                text_parts.extend(slide_texts)

        if not text_parts:
            return "Kein Text in der PowerPoint-Datei gefunden."

        return "\n".join(text_parts)

    except Exception as e:
        return f"Fehler beim Lesen der PowerPoint-Datei: {e}"
